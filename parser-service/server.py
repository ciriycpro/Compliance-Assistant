"""
Parser Service v1 — по DEC-008.

Детерминированный роутер по MIME → L1-L14 библиотеки.
Уровень A: switch по MIME (~80% случаев).
Уровень B: PDF-ветка с pdf-inspector классификацией + fallback цепочка.

L14 (LLM-vision Qwen3-VL + Qwen 2.5 VL fallback) — только для растровых форматов.

Endpoint: POST /parse с {path: "/var/lib/mail-stack/attachments/.../file.ext"}
Возврат: единый JSON формат по DEC-008.
"""

import os
import io
import csv
import json
import logging
import mimetypes
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional

import httpx
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel

# === Конфиг через env (docker-friendly по DEC-007) ===
OPENROUTER_API_KEY = os.getenv('OPENROUTER_API_KEY')
OPENROUTER_VISION_PRIMARY = os.getenv('OPENROUTER_VISION_PRIMARY', 'qwen/qwen3-vl-235b-a22b-instruct')
OPENROUTER_VISION_FALLBACK = os.getenv('OPENROUTER_VISION_FALLBACK', 'qwen/qwen2.5-vl-72b-instruct')
OPENROUTER_TIMEOUT_SEC = int(os.getenv('OPENROUTER_TIMEOUT_SEC', '120'))
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

LLM_VISION_PROMPT = os.getenv(
    'LLM_VISION_PROMPT',
    'Извлеки весь текст из этого документа на русском языке. Если есть таблицы — '
    'сохрани структуру через markdown-таблицы. Если есть печати — отметь их в тексте '
    'в формате [ПЕЧАТЬ: <текст печати>]. Не добавляй пояснений и комментариев, только текст документа.'
)

PDF_CONFIDENCE_THRESHOLD = float(os.getenv('PDF_CONFIDENCE_THRESHOLD', '0.8'))
MIN_EXTRACTED_TEXT_LEN = int(os.getenv('MIN_EXTRACTED_TEXT_LEN', '50'))

# === Логирование в stdout ===
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
)
log = logging.getLogger("parser-service")


# === Модели request/response ===
class ParseRequest(BaseModel):
    path: str


class ParseResponse(BaseModel):
    text: str
    method: str
    format: str
    pages: Optional[int] = None
    warnings: list[str] = []
    cost_estimate_usd: Optional[float] = None


# ========================================================================
# === ПАРСЕРЫ — каждая функция возвращает (text, pages_or_none, warnings) ===
# ========================================================================


# === L5: DOCX (mammoth) ===
def parse_docx(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    import mammoth
    warnings = []
    with io.BytesIO(file_bytes) as buf:
        result = mammoth.extract_raw_text(buf)
        text = result.value or ""
        if result.messages:
            warnings.extend([m.message for m in result.messages[:5]])
    return text.strip(), None, warnings


# === L6: XLSX (openpyxl) ===
# vNext: для банковских выписок — семантический фильтр по колонке "Назначение платежа"
# (формат разный у Сбера/ТКБ/Совкомбанка, требует отдельной логики на Compliance-этапе)
XLSX_ROW_HARD_LIMIT = int(os.getenv('XLSX_ROW_HARD_LIMIT', '200'))
XLSX_KEEP_FIRST = int(os.getenv('XLSX_KEEP_FIRST', '100'))
XLSX_KEEP_LAST = int(os.getenv('XLSX_KEEP_LAST', '50'))


def parse_xlsx(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    from openpyxl import load_workbook
    warnings = []
    with io.BytesIO(file_bytes) as buf:
        wb = load_workbook(buf, data_only=True, read_only=True)
        chunks = []
        for ws in wb.worksheets:
            chunks.append(f"=== Лист: {ws.title} ===")

            # Собираем все непустые строки
            all_rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join('' if v is None else str(v) for v in row)
                if row_text.strip():
                    all_rows.append(row_text)

            if len(all_rows) <= XLSX_ROW_HARD_LIMIT:
                chunks.extend(all_rows)
            else:
                chunks.extend(all_rows[:XLSX_KEEP_FIRST])
                cut = len(all_rows) - XLSX_KEEP_FIRST - XLSX_KEEP_LAST
                chunks.append(
                    f"\n[ВЫРЕЗАНО {cut} строк из середины. "
                    f"Всего строк: {len(all_rows)}. Полный файл — в исходном письме.]\n"
                )
                chunks.extend(all_rows[-XLSX_KEEP_LAST:])
                warnings.append(f"xlsx_truncated:sheet={ws.title}:total={len(all_rows)}")

        wb.close()
    return '\n'.join(chunks).strip(), None, warnings


# === L7: XLS старый (xlrd) ===
def parse_xls(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    import xlrd
    warnings = []
    wb = xlrd.open_workbook(file_contents=file_bytes)
    chunks = []
    for sheet in wb.sheets():
        chunks.append(f"=== Лист: {sheet.name} ===")
        for row_idx in range(sheet.nrows):
            row = sheet.row_values(row_idx)
            row_text = '\t'.join('' if v == '' else str(v) for v in row)
            if row_text.strip():
                chunks.append(row_text)
    return '\n'.join(chunks).strip(), None, warnings


# === L8: PPTX (python-pptx) ===
def parse_pptx(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    from pptx import Presentation
    warnings = []
    with io.BytesIO(file_bytes) as buf:
        prs = Presentation(buf)
        chunks = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks.append(f"=== Слайд {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = ''.join(run.text for run in para.runs)
                        if para_text.strip():
                            chunks.append(para_text)
    return '\n'.join(chunks).strip(), len(prs.slides), warnings


# === L9: HTML (beautifulsoup4) ===
def parse_html(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    from bs4 import BeautifulSoup
    warnings = []
    try:
        text_bytes = file_bytes.decode('utf-8', errors='replace')
    except Exception:
        text_bytes = file_bytes.decode('cp1251', errors='replace')
    soup = BeautifulSoup(text_bytes, 'lxml')
    # Удаляем скрипты и стили
    for s in soup(['script', 'style']):
        s.decompose()
    text = soup.get_text(separator='\n', strip=True)
    return text, None, warnings


# === L10: CSV ===
def parse_csv(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    warnings = []
    text_bytes = file_bytes.decode('utf-8', errors='replace')
    chunks = []
    reader = csv.reader(io.StringIO(text_bytes))
    for row in reader:
        chunks.append('\t'.join(row))
    return '\n'.join(chunks).strip(), None, warnings


# === L11: TXT ===
def parse_txt(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    warnings = []
    try:
        text = file_bytes.decode('utf-8', errors='replace')
    except Exception:
        text = file_bytes.decode('cp1251', errors='replace')
    return text.strip(), None, warnings


# === L12: XML ===
def parse_xml(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    warnings = []
    text_bytes = file_bytes.decode('utf-8', errors='replace')
    try:
        root = ET.fromstring(text_bytes)
        chunks = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                chunks.append(f"{tag}: {elem.text.strip()}")
        return '\n'.join(chunks).strip(), None, warnings
    except ET.ParseError as e:
        warnings.append(f"xml_parse_error: {e}")
        return text_bytes, None, warnings


# === L13: JSON ===
def parse_json_file(file_bytes: bytes) -> tuple[str, Optional[int], list[str]]:
    warnings = []
    text_bytes = file_bytes.decode('utf-8', errors='replace')
    try:
        data = json.loads(text_bytes)
        text = json.dumps(data, ensure_ascii=False, indent=2)
        return text, None, warnings
    except json.JSONDecodeError as e:
        warnings.append(f"json_parse_error: {e}")
        return text_bytes, None, warnings


# === L3: PyMuPDF извлечение PDF ===
def extract_pdf_pymupdf(file_bytes: bytes) -> tuple[str, int, list[str]]:
    import fitz
    warnings = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks = []
    for page in doc:
        text = page.get_text()
        if text:
            chunks.append(text)
    pages = len(doc)
    doc.close()
    return '\n'.join(chunks).strip(), pages, warnings


# === L4: pdfplumber извлечение PDF (fallback) ===
def extract_pdf_pdfplumber(file_bytes: bytes) -> tuple[str, int, list[str]]:
    import pdfplumber
    warnings = []
    chunks = []
    with io.BytesIO(file_bytes) as buf:
        with pdfplumber.open(buf) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text:
                    chunks.append(text)
            pages = len(pdf.pages)
    return '\n'.join(chunks).strip(), pages, warnings


# === L1: pdf-inspector классификация PDF ===
def classify_pdf(file_bytes: bytes) -> tuple[str, float, list[int]]:
    """
    Возвращает (type, confidence, pages_needing_ocr).
    type ∈ {"TextBased", "Scanned", "ImageBased", "Mixed", "Unknown"}

    Использует pdf_inspector.detect_pdf_bytes() — fast detection без извлечения текста.
    pdf_inspector возвращает pdf_type как lowercase ('text_based', 'scanned', ...),
    нормализуем к camelCase для совместимости с DEC-008.
    """
    try:
        import pdf_inspector
        result = pdf_inspector.detect_pdf_bytes(file_bytes)
        # Нормализация имён типов из snake_case в CamelCase
        type_map = {
            'text_based': 'TextBased',
            'scanned': 'Scanned',
            'image_based': 'ImageBased',
            'mixed': 'Mixed',
        }
        pdf_type = type_map.get(result.pdf_type, result.pdf_type)
        return (
            pdf_type,
            result.confidence,
            result.pages_needing_ocr,
        )
    except (ImportError, AttributeError) as e:
        log.warning(f"pdf-inspector unavailable ({e}), fallback to PyMuPDF heuristic")
        return classify_pdf_pymupdf_fallback(file_bytes)
    except Exception as e:
        log.warning(f"pdf-inspector classification failed ({e}), fallback to PyMuPDF")
        return classify_pdf_pymupdf_fallback(file_bytes)


def classify_pdf_pymupdf_fallback(file_bytes: bytes) -> tuple[str, float, list[int]]:
    """L2 (fallback): классификация через PyMuPDF — считаем text vs image на страницах."""
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_with_text = 0
    pages_with_image_only = 0
    pages_needing_ocr = []
    total = len(doc)

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        has_text = len(text) > MIN_EXTRACTED_TEXT_LEN
        has_image = bool(page.get_images())
        if has_text:
            pages_with_text += 1
        elif has_image:
            pages_with_image_only += 1
            pages_needing_ocr.append(i)
    doc.close()

    if pages_with_text == total:
        return "TextBased", 1.0, []
    if pages_with_image_only > 0 and pages_with_text == 0:
        return "Scanned", 1.0, pages_needing_ocr
    if pages_with_text > 0 and pages_with_image_only > 0:
        return "Mixed", 0.7, pages_needing_ocr
    return "Unknown", 0.0, []


# === L14: LLM-vision (Qwen3-VL + fallback) ===
import base64


def vision_parse_image(file_bytes: bytes, mime: str = "image/jpeg") -> tuple[str, str, float, list[str]]:
    """
    Отправляет изображение в OpenRouter Qwen3-VL.
    Возвращает (text, model_used, cost_usd, warnings).
    """
    if not OPENROUTER_API_KEY:
        raise HTTPException(
            status_code=500,
            detail="OPENROUTER_API_KEY not configured for L14 vision",
        )

    b64 = base64.b64encode(file_bytes).decode('ascii')
    data_url = f"data:{mime};base64,{b64}"

    payload = {
        "model": OPENROUTER_VISION_PRIMARY,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": LLM_VISION_PROMPT},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ],
    }

    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }

    warnings = []

    # Primary: Qwen3-VL 235B
    try:
        with httpx.Client(timeout=OPENROUTER_TIMEOUT_SEC) as client:
            r = client.post(OPENROUTER_URL, json=payload, headers=headers)
            if r.status_code == 200:
                data = r.json()
                text = data["choices"][0]["message"]["content"]
                usage = data.get("usage", {})
                cost = estimate_cost_qwen3vl(usage)
                return text.strip(), OPENROUTER_VISION_PRIMARY, cost, warnings
            warnings.append(f"primary_failed_{r.status_code}")
            log.warning(f"Vision primary failed: {r.status_code} {r.text[:200]}")
    except Exception as e:
        warnings.append(f"primary_exception: {type(e).__name__}")
        log.warning(f"Vision primary exception: {e}")

    # Fallback: Qwen 2.5 VL 72B
    payload["model"] = OPENROUTER_VISION_FALLBACK
    try:
        with httpx.Client(timeout=OPENROUTER_TIMEOUT_SEC) as client:
            r = client.post(OPENROUTER_URL, json=payload, headers=headers)
            if r.status_code == 200:
                data = r.json()
                text = data["choices"][0]["message"]["content"]
                usage = data.get("usage", {})
                cost = estimate_cost_qwen25vl(usage)
                warnings.append("fallback_used")
                return text.strip(), OPENROUTER_VISION_FALLBACK, cost, warnings
            warnings.append(f"fallback_failed_{r.status_code}")
            log.error(f"Vision fallback failed: {r.status_code} {r.text[:200]}")
    except Exception as e:
        warnings.append(f"fallback_exception: {type(e).__name__}")
        log.error(f"Vision fallback exception: {e}")

    raise HTTPException(
        status_code=502,
        detail={"error": "vision_all_providers_failed", "warnings": warnings},
    )


def estimate_cost_qwen3vl(usage: dict) -> float:
    """Qwen3-VL 235B: $0.20/M input, $0.88/M output."""
    in_tok = usage.get("prompt_tokens", 0)
    out_tok = usage.get("completion_tokens", 0)
    return round(in_tok * 0.20e-6 + out_tok * 0.88e-6, 5)


def estimate_cost_qwen25vl(usage: dict) -> float:
    """Qwen 2.5 VL 72B: ~$0.15/M input/output (DeepInfra)."""
    in_tok = usage.get("prompt_tokens", 0)
    out_tok = usage.get("completion_tokens", 0)
    return round((in_tok + out_tok) * 0.15e-6, 5)


def vision_parse_pdf(file_bytes: bytes) -> tuple[str, str, float, list[str]]:
    """Рендерит каждую страницу PDF в PNG и отправляет в LLM-vision."""
    import fitz
    warnings = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks = []
    total_cost = 0.0
    model_used = None

    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        png_bytes = pix.tobytes("png")
        try:
            text, model, cost, wrs = vision_parse_image(png_bytes, "image/png")
            chunks.append(f"=== Страница {i+1} ===\n{text}")
            total_cost += cost
            model_used = model
            warnings.extend(wrs)
        except HTTPException as e:
            warnings.append(f"page_{i+1}_failed")
            chunks.append(f"=== Страница {i+1} ===\n[OCR FAILED]")

    pages = len(doc)
    doc.close()
    return '\n\n'.join(chunks).strip(), model_used or "vision_failed", total_cost, warnings


# ========================================================================
# === MAIN ROUTER (Уровень A + Уровень B) ===
# ========================================================================


def detect_mime_from_path(path: Path) -> str:
    """Определяем MIME по расширению."""
    mime, _ = mimetypes.guess_type(str(path))
    if mime:
        return mime
    # Дополнительные расширения которых нет в стандартной mimetypes
    ext = path.suffix.lower()
    return {
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.xls': 'application/vnd.ms-excel',
        '.heic': 'image/heic',
    }.get(ext, 'application/octet-stream')


def parse_pdf_mixed_per_page(file_bytes: bytes, pages_needing_ocr: list[int]) -> tuple[str, str, float, list[str]]:
    """
    Mixed PDF: per-page роутинг (DEC-008 уровень B).

    ИЗМЕНЕНО 2026-05-14: inspector подсказывает, но не диктует. Сначала пробуем
    text-extraction на КАЖДОЙ странице — даже на тех что inspector пометил OCR.
    Vision только если текстовый слой реально пуст/слишком короткий.

    Лечит ложные срабатывания pdf_inspector на формах с штрихкодами/QR-кодами
    (налоговые декларации, банковские формы и т.п.), где текстовый слой есть,
    но inspector видит "графический шум" и думает что нужен OCR.
    """
    import fitz
    warnings = ["per_page_routing_text_first"]
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks = []
    total_cost = 0.0
    model_used = "pymupdf"

    ocr_hint_set = set(pages_needing_ocr)
    text_pages = 0
    vision_pages = 0
    rescued_pages = 0  # inspector пометил OCR, но text справился

    for i, page in enumerate(doc):
        page_num = i + 1
        text = page.get_text().strip()

        if text and len(text) >= MIN_EXTRACTED_TEXT_LEN:
            chunks.append(f"=== Страница {page_num} (text) ===\n{text}")
            text_pages += 1
            if i in ocr_hint_set:
                rescued_pages += 1
        else:
            try:
                pix = page.get_pixmap(dpi=150)
                png_bytes = pix.tobytes("png")
                v_text, v_model, v_cost, v_wrs = vision_parse_image(png_bytes, "image/png")
                chunks.append(f"=== Страница {page_num} (vision) ===\n{v_text}")
                total_cost += v_cost
                model_used = v_model
                vision_pages += 1
                warnings.extend(v_wrs)
            except HTTPException as e:
                warnings.append(f"page_{page_num}_vision_failed")
                chunks.append(f"=== Страница {page_num} ===\n[OCR FAILED]")

    doc.close()

    warnings.append(f"pages_text:{text_pages}")
    warnings.append(f"pages_vision:{vision_pages}")
    if rescued_pages:
        warnings.append(f"pages_rescued_from_ocr:{rescued_pages}")

    return '\n\n'.join(chunks).strip(), model_used, total_cost, warnings


def parse_pdf_with_routing(file_bytes: bytes) -> tuple[str, str, str, int, list[str], float]:
    """
    Уровень B: маршрутизация внутри PDF.
    Возвращает (text, method, format, pages, warnings, cost).
    """
    pdf_type, confidence, pages_needing_ocr = classify_pdf(file_bytes)
    log.info(f"PDF classification: type={pdf_type}, conf={confidence}, ocr_pages={len(pages_needing_ocr)}")
    warnings = [f"pdf_class:{pdf_type}", f"conf:{confidence}"]

    if pdf_type == "TextBased" and confidence > PDF_CONFIDENCE_THRESHOLD:
        # L3 primary
        try:
            text, pages, wrs = extract_pdf_pymupdf(file_bytes)
            warnings.extend(wrs)
            if text and len(text) >= MIN_EXTRACTED_TEXT_LEN:
                return text, "pymupdf", "pdf-text", pages, warnings, 0.0
            warnings.append("pymupdf_too_short")
        except Exception as e:
            warnings.append(f"pymupdf_error:{type(e).__name__}")
            log.warning(f"PyMuPDF extract failed: {e}")

        # L4 fallback
        try:
            text, pages, wrs = extract_pdf_pdfplumber(file_bytes)
            warnings.extend(wrs)
            if text and len(text) >= MIN_EXTRACTED_TEXT_LEN:
                warnings.append("pdfplumber_fallback_used")
                return text, "pdfplumber", "pdf-text", pages, warnings, 0.0
            warnings.append("pdfplumber_too_short")
        except Exception as e:
            warnings.append(f"pdfplumber_error:{type(e).__name__}")
            log.warning(f"pdfplumber extract failed: {e}")

        # L14 last resort
        warnings.append("text_extraction_failed_fallback_to_vision")
        text, model, cost, wrs = vision_parse_pdf(file_bytes)
        warnings.extend(wrs)
        return text, model, "pdf-vision-fallback", None, warnings, cost

    if pdf_type in ("Scanned", "ImageBased"):
        # L14 сразу
        text, model, cost, wrs = vision_parse_pdf(file_bytes)
        warnings.extend(wrs)
        return text, model, "pdf-scan", None, warnings, cost

    # Mixed / Unknown — честный per-page роутинг (DEC-008 уровень B)
    # PyMuPDF для страниц с текстом, L14 (LLM-vision) только для image-only страниц
    if pdf_type == "Unknown":
        # Unknown — на всякий случай прогоняем PyMuPDF, считаем недостающие страницы
        _, _, pages_needing_ocr = classify_pdf_pymupdf_fallback(file_bytes)
    text, model, cost, wrs = parse_pdf_mixed_per_page(file_bytes, pages_needing_ocr)
    warnings.extend(wrs)
    fmt = "pdf-mixed-per-page" if model.startswith("qwen") else "pdf-mixed-text-only"
    return text, model, fmt, None, warnings, cost


# ========================================================================
# === FastAPI app ===
# ========================================================================
app = FastAPI(title="parser-service-v1")


@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "parser-service-v1",
        "openrouter_configured": bool(OPENROUTER_API_KEY),
        "vision_primary": OPENROUTER_VISION_PRIMARY,
        "vision_fallback": OPENROUTER_VISION_FALLBACK,
    }


@app.post("/parse", response_model=ParseResponse)
def parse(req: ParseRequest):
    """
    Парсинг файла по пути.
    Алгоритм (DEC-008):
    1. Уровень A: switch по MIME → специализированный парсер
    2. Если MIME=application/pdf → Уровень B (PDF-ветка с классификацией)
    3. Если image/* → L14 LLM-vision
    """
    path = Path(req.path)
    if not path.exists() or not path.is_file():
        raise HTTPException(status_code=404, detail=f"File not found: {req.path}")

    mime = detect_mime_from_path(path)
    log.info(f"Parsing {path} (mime={mime}, size={path.stat().st_size})")

    with open(path, 'rb') as f:
        file_bytes = f.read()

    try:
        # === Уровень A: switch по MIME ===
        if mime == 'application/pdf':
            text, method, fmt, pages, warnings, cost = parse_pdf_with_routing(file_bytes)
            return ParseResponse(
                text=text, method=method, format=fmt,
                pages=pages, warnings=warnings, cost_estimate_usd=cost or None,
            )

        if mime.startswith('image/'):
            text, model, cost, warnings = vision_parse_image(file_bytes, mime)
            return ParseResponse(
                text=text, method=model, format=mime.split('/')[-1],
                pages=None, warnings=warnings, cost_estimate_usd=cost or None,
            )

        if mime.endswith('wordprocessingml.document'):
            text, pages, wrs = parse_docx(file_bytes)
            return ParseResponse(text=text, method="mammoth", format="docx", pages=pages, warnings=wrs)

        if mime.endswith('spreadsheetml.sheet'):
            text, pages, wrs = parse_xlsx(file_bytes)
            return ParseResponse(text=text, method="openpyxl", format="xlsx", pages=pages, warnings=wrs)

        if mime == 'application/vnd.ms-excel':
            text, pages, wrs = parse_xls(file_bytes)
            return ParseResponse(text=text, method="xlrd", format="xls", pages=pages, warnings=wrs)

        if mime.endswith('presentationml.presentation'):
            text, pages, wrs = parse_pptx(file_bytes)
            return ParseResponse(text=text, method="python-pptx", format="pptx", pages=pages, warnings=wrs)

        if mime in ('text/html', 'application/xhtml+xml'):
            text, pages, wrs = parse_html(file_bytes)
            return ParseResponse(text=text, method="beautifulsoup4", format="html", pages=pages, warnings=wrs)

        if mime == 'text/csv':
            text, pages, wrs = parse_csv(file_bytes)
            return ParseResponse(text=text, method="csv", format="csv", pages=pages, warnings=wrs)

        if mime in ('text/plain',):
            text, pages, wrs = parse_txt(file_bytes)
            return ParseResponse(text=text, method="read", format="txt", pages=pages, warnings=wrs)

        if mime in ('application/xml', 'text/xml'):
            text, pages, wrs = parse_xml(file_bytes)
            return ParseResponse(text=text, method="xml.etree", format="xml", pages=pages, warnings=wrs)

        if mime == 'application/json':
            text, pages, wrs = parse_json_file(file_bytes)
            return ParseResponse(text=text, method="json", format="json", pages=pages, warnings=wrs)

        # Неизвестный MIME
        raise HTTPException(
            status_code=415,
            detail={"error": "unsupported_mime", "mime": mime, "path": str(path)},
        )

    except HTTPException:
        raise
    except Exception as e:
        log.error(f"Parse failed for {path}: {e}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail={"error": "parser_exception", "message": str(e), "type": type(e).__name__},
        )
